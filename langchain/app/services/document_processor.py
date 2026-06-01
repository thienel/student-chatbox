import asyncio
import logging
import os
import uuid
from pathlib import Path

import httpx
from langchain_openai import OpenAIEmbeddings
from langchain_text_splitters import RecursiveCharacterTextSplitter

from ..config import settings
from .qdrant_service import qdrant_service

logger = logging.getLogger(__name__)

_splitter = RecursiveCharacterTextSplitter(chunk_size=1000, chunk_overlap=200)


def _get_embeddings() -> OpenAIEmbeddings:
    kwargs: dict = {
        "openai_api_key": settings.openai_api_key,
        "model": settings.openai_embedding_model,
        "chunk_size": 100,
    }
    if settings.openai_base_url:
        kwargs["base_url"] = settings.openai_base_url
    return OpenAIEmbeddings(**kwargs)


def _extract_text(file_path: str) -> str:
    ext = Path(file_path).suffix.lower()
    if ext == ".pdf":
        return _extract_pdf(file_path)
    if ext == ".docx":
        return _extract_docx(file_path)
    if ext == ".pptx":
        return _extract_pptx(file_path)
    raise ValueError(f"Unsupported file type: {ext}")


def _extract_pdf(file_path: str) -> str:
    import fitz  # type: ignore[import]

    doc = fitz.open(file_path)
    return "\n".join(page.get_text() for page in doc)


def _extract_docx(file_path: str) -> str:
    from docx import Document  # type: ignore[import]

    return "\n".join(p.text for p in Document(file_path).paragraphs)


def _extract_pptx(file_path: str) -> str:
    from pptx import Presentation  # type: ignore[import]

    prs = Presentation(file_path)
    texts = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                texts.append(shape.text)
    return "\n".join(texts)


async def _notify_nestjs(document_id: str, status: str, chunk_count: int = 0, error: str = "") -> None:
    url = f"{settings.nestjs_url}/internal/documents/{document_id}/processing-result"
    payload = {"status": status, "chunkCount": chunk_count, "error": error}
    try:
        async with httpx.AsyncClient(timeout=10) as client:
            await client.patch(
                url,
                json=payload,
                headers={"x-internal-key": settings.ai_service_secret},
            )
    except Exception as exc:
        logger.error(f"Failed to notify NestJS for document {document_id}: {exc}")


async def process_document(document_id: str, file_path: str, subject_id: str) -> None:
    try:
        logger.info(f"Processing document {document_id} at {file_path}")

        text = await asyncio.to_thread(_extract_text, file_path)
        if not text or not text.strip():
            raise ValueError("Could not extract text from document")

        chunks = _splitter.split_text(text)
        non_empty = [c for c in chunks if c.strip()]
        if not non_empty:
            raise ValueError("Document produced no chunks after splitting")

        logger.info(f"Embedding {len(non_empty)} chunks…")
        embeddings = _get_embeddings()
        vectors: list[list[float]] = await asyncio.to_thread(embeddings.embed_documents, non_empty)

        vector_dim = len(vectors[0])
        logger.info(f"Embedded {len(vectors)} vectors (dim={vector_dim}), upserting to Qdrant…")

        original_name = os.path.basename(file_path)
        # Strip timestamp prefix added by NestJS LocalFileService (e.g. "1234567890_name.pdf")
        if "_" in original_name:
            original_name = "_".join(original_name.split("_")[1:])

        await asyncio.to_thread(qdrant_service.ensure_collection, vector_dim)

        points = [
            {
                "id": str(uuid.uuid4()),
                "vector": vectors[i],
                "payload": {
                    "document_id": document_id,
                    "subject_id": subject_id,
                    "chunk_index": i,
                    "text": chunk,
                    "original_name": original_name,
                },
            }
            for i, chunk in enumerate(non_empty)
        ]

        await asyncio.to_thread(qdrant_service.upsert_points, points)
        logger.info(f"Document {document_id} processed: {len(non_empty)} chunks")

        await _notify_nestjs(document_id, "ready", len(non_empty))

    except Exception as exc:
        logger.error(f"Document {document_id} processing failed: {exc}")
        await _notify_nestjs(document_id, "failed", 0, str(exc))
